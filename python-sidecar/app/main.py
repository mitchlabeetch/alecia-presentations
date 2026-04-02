"""
Alecia Presentation Builder - Python Sidecar
Service de conversion de documents PPTX/PDF via Docling
"""

import os
import json
import logging
import tempfile
import asyncio
from typing import Optional, Any
from pathlib import Path
from contextlib import asynccontextmanager

from fastapi import FastAPI, File, UploadFile, HTTPException, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import aiofiles

# Configuration du logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger("alecia-docling")

# ============================================================================
# Modèles de données
# ============================================================================

class SlideContent(BaseModel):
    """Modèle de contenu pour une slide Alecia"""
    text: Optional[str] = None
    subtitle: Optional[str] = None
    items: Optional[list[str]] = None
    image_url: Optional[str] = None
    table_data: Optional[dict] = None
    chart_type: Optional[str] = None
    chart_data: Optional[dict] = None
    quote: Optional[str] = None
    author: Optional[str] = None
    kpis: Optional[list[dict]] = None
    timeline: Optional[list[dict]] = None
    team: Optional[list[dict]] = None
    sections: Optional[list[str]] = None
    metadata: Optional[dict] = None

class SlideResult(BaseModel):
    """Modèle de résultat pour une slide"""
    order_index: int
    type: str
    title: str
    content: SlideContent
    notes: Optional[str] = None
    image_path: Optional[str] = None

class ConversionResponse(BaseModel):
    """Réponse standard pour la conversion"""
    status: str
    slides: Optional[list[SlideResult]] = None
    document: Optional[dict] = None
    error: Optional[str] = None
    filename: Optional[str] = None
    slide_count: Optional[int] = None

class BatchItemResult(BaseModel):
    """Résultat pour un fichier dans un batch"""
    filename: str
    status: str
    slides: Optional[list[SlideResult]] = None
    error: Optional[str] = None

class BatchResponse(BaseModel):
    """Réponse pour conversion batch"""
    status: str
    results: list[BatchItemResult]
    total_count: int
    success_count: int
    error_count: int

class JobStatus(BaseModel):
    """Statut d'un job de conversion"""
    job_id: str
    status: str
    filename: Optional[str] = None
    progress: int = 0
    result: Optional[ConversionResponse] = None

# ============================================================================
# Gestionnaire de cycle de vie
# ============================================================================

@asynccontextmanager
async def lifespan(app: FastAPI):
    """Gestionnaire de cycle de vie pour démarrage/arrêt"""
    logger.info("🚀 Alecia Docling Sidecar - Démarrage du service...")
    logger.info("📄 Services disponibles: PPTX, PDF")
    logger.info("🔗 CORS activé pour: http://localhost:3000, http://localhost:3001")

    # Initialisation si nécessaire
    app.state.jobs = {}

    yield

    logger.info("👋 Alecia Docling Sidecar - Arrêt du service...")

# ============================================================================
# Application FastAPI
# ============================================================================

app = FastAPI(
    title="Alecia Docling Import Service",
    description="""
    Service de conversion de documents pour Alecia Presentation Builder.

    ## Fonctionnalités
    - Conversion PPTX en structure Alecia
    - Conversion PDF en structure Alecia
    - Conversion par lot (batch)
    - Suivi de progression des jobs

    ## Formats supportés
    - PowerPoint: .pptx, .ppt
    - PDF: .pdf
    """,
    version="1.0.0",
    lifespan=lifespan,
)

# Middleware CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://localhost:3001",
        "http://127.0.0.1:3000",
        "http://127.0.0.1:3001"
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ============================================================================
# Routes de santé
# ============================================================================

@app.get("/health")
async def health_check():
    """Vérification de l'état du service"""
    return {
        "status": "ok",
        "service": "alecia-docling",
        "version": "1.0.0"
    }

@app.get("/")
async def root():
    """Route racine avec informations du service"""
    return {
        "service": "Alecia Docling Import Service",
        "version": "1.0.0",
        "endpoints": {
            "health": "/health",
            "convert_pptx": "POST /convert/pptx",
            "convert_pdf": "POST /convert/pdf",
            "convert_batch": "POST /convert/batch",
            "job_status": "GET /status/{job_id}",
            "validate": "POST /validate"
        }
    }

# ============================================================================
# Fonctions utilitaires de conversion
# ============================================================================

def extract_text_content(element: Any) -> str:
    """Extrait le texte d'un élément Docling"""
    if hasattr(element, 'text'):
        return element.text
    if hasattr(element, 'content'):
        return str(element.content)
    return ""

def infer_slide_type(title: str, content: str) -> str:
    """Inférence le type de slide basé sur le contenu"""

    title_lower = title.lower()
    content_lower = content.lower()

    # Slides especiales
    if 'sommaire' in title_lower or 'summary' in title_lower:
        return 'Section_Navigator'
    if 'couverture' in title_lower or 'cover' in title_lower:
        return 'Cover'
    if 'contact' in title_lower:
        return 'Contact_Block'

    # Slides de contenu
    if any(kw in title_lower for kw in ['swot', 'forces', 'faiblesses']):
        return 'SWOT'
    if any(kw in title_lower for kw in ['kpi', 'indicateur', 'métrique']):
        return 'KPI_Card'
    if any(kw in title_lower for kw in ['tableau', 'table', 'comparaison']):
        return 'Table_Block'
    if any(kw in title_lower for kw in ['graphique', 'chart', 'graph']):
        return 'Chart_Block'
    if any(kw in title_lower for kw in ['timeline', 'chronolog', 'jalon']):
        return 'Process_Timeline'
    if any(kw in title_lower for kw in ['citation', 'quote']):
        return 'Quote_Block'
    if any(kw in title_lower for kw in ['équipe', 'team', 'membre']):
        return 'Team_Members'
    if any(kw in title_lower for kw in ['entreprise', 'société', 'company']):
        return 'Company_Overview'
    if any(kw in title_lower for kw in ['track record', 'références']):
        return 'Trackrecord_Block'
    if any(kw in title_lower for kw in ['rationale', 'rational']):
        return 'Deal_Rationale'

    return 'Text_Block'

def parse_pptx_slides(pptx_path: str) -> list[SlideResult]:
    """
    Parse un fichier PPTX et extrait les slides.

    Utilise python-pptx pour l'extraction basique.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    slides = []

    try:
        prs = Presentation(pptx_path)

        for idx, slide in enumerate(prs.slides):
            title = ""
            content_text = []
            notes_text = None
            image_paths = []

            # Extraire le titre
            if slide.shapes.title:
                title = slide.shapes.title.text or ""

            # Extraire le contenu
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape != slide.shapes.title:
                    text = shape.text.strip()
                    if text:
                        content_text.append(text)

                # Détecter les images
                if hasattr(shape, 'image'):
                    try:
                        image = shape.image
                        image_paths.append(f"image_{idx}_{len(image_paths)}.png")
                    except:
                        pass

            # Extraire les notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text

            # Construire le contenu
            full_content = "\n".join(content_text)

            # Créer le SlideResult
            slide_type = infer_slide_type(title, full_content)

            slide_result = SlideResult(
                order_index=idx,
                type=slide_type,
                title=title or f"Slide {idx + 1}",
                content=SlideContent(
                    text=full_content if content_text else None,
                    subtitle=None,
                    items=content_text if slide_type == 'Section_Navigator' else None,
                ),
                notes=notes_text,
                image_path=image_paths[0] if image_paths else None
            )

            slides.append(slide_result)

    except ImportError:
        logger.warning("python-pptx non disponible, utilisation du fallback")
        # Fallback basique
        slides.append(SlideResult(
            order_index=0,
            type='Cover',
            title='Import PPTX',
            content=SlideContent(text='Contenu importé depuis PowerPoint'),
            notes=None,
            image_path=None
        ))

    return slides

def parse_pdf_pages(pdf_path: str) -> list[SlideResult]:
    """
    Parse un fichier PDF et extrait les pages comme slides.

    Utilise une approche basique pour extraire le texte.
    """
    slides = []

    try:
        import fitz  # PyMuPDF

        doc = fitz.open(pdf_path)

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()

            # Parser le texte pour extraire titre et contenu
            lines = text.strip().split('\n')
            title = lines[0] if lines else f"Page {page_num + 1}"
            content = '\n'.join(lines[1:]) if len(lines) > 1 else text

            slide_type = infer_slide_type(title, content)

            slide_result = SlideResult(
                order_index=page_num,
                type=slide_type,
                title=title.strip(),
                content=SlideContent(text=content.strip()),
                notes=None,
                image_path=None
            )

            slides.append(slide_result)

        doc.close()

    except ImportError:
        logger.warning("PyMuPDF non disponible")
        slides.append(SlideResult(
            order_index=0,
            type='Cover',
            title='Import PDF',
            content=SlideContent(text='Contenu importé depuis PDF'),
            notes=None,
            image_path=None
        ))

    return slides

# ============================================================================
# Routes de conversion
# ============================================================================

@app.post("/convert/pptx", response_model=ConversionResponse)
async def convert_pptx(
    file: UploadFile = File(...),
    background_tasks: BackgroundTasks = None
):
    """
    Convertit un fichier PPTX en structure Alecia.

    - **file**: Fichier PowerPoint (.pptx, .ppt)

    Retourne une liste de slides structurées.
    """
    # Valider le type de fichier
    filename = file.filename or "unnamed.pptx"
    if not filename.lower().endswith(('.pptx', '.ppt')):
        raise HTTPException(
            status_code=400,
            detail="Seuls les fichiers PPTX et PPT sont supportés"
        )

    logger.info(f"📄 Conversion PPTX: {filename}")

    # Sauvegarder temporairement le fichier
    with tempfile.NamedTemporaryFile(
        delete=False,
        suffix=os.path.splitext(filename)[1]
    ) as tmp_file:
        content = await file.read()
        tmp_file.write(content)
        tmp_path = tmp_file.name

    try:
        # Parser le PPTX
        slides = parse_pptx_slides(tmp_path)

        logger.info(f"✅ {len(slides)} slides extraites de {filename}")

        return ConversionResponse(
            status="success",
            slides=slides,
            filename=filename,
            slide_count=len(slides)
        )

    except Exception as e:
        logger.error(f"❌ Erreur de conversion PPTX: {str(e)}")
        return ConversionResponse(
            status="error",
            error=str(e),
            filename=filename
        )

    finally:
        # Nettoyage
        try:
            os.unlink(tmp_path)
        except:
            pass


@app.post("/convert/pdf", response_model=ConversionResponse)
async def convert_pdf(
    file: UploadFile = File(...),
    background_tasks: BackgroundTasks = None
):
    """
    Convertit un fichier PDF en structure Alecia.

    - **file**: Fichier PDF (.pdf)

    Retourne une liste de slides structurées (une par page).
    """
    # Valider le type de fichier
    filename = file.filename or "unnamed.pdf"
    if not filename.lower().endswith('.pdf'):
        raise HTTPException(
            status_code=400,
            detail="Seuls les fichiers PDF sont supportés"
        )

    logger.info(f"📄 Conversion PDF: {filename}")

    # Sauvegarder temporairement le fichier
    with tempfile.NamedTemporaryFile(
        delete=False,
        suffix=".pdf"
    ) as tmp_file:
        content = await file.read()
        tmp_file.write(content)
        tmp_path = tmp_file.name

    try:
        # Parser le PDF
        slides = parse_pdf_pages(tmp_path)

        logger.info(f"✅ {len(slides)} slides extraites de {filename}")

        return ConversionResponse(
            status="success",
            slides=slides,
            filename=filename,
            slide_count=len(slides)
        )

    except Exception as e:
        logger.error(f"❌ Erreur de conversion PDF: {str(e)}")
        return ConversionResponse(
            status="error",
            error=str(e),
            filename=filename
        )

    finally:
        # Nettoyage
        try:
            os.unlink(tmp_path)
        except:
            pass


@app.post("/convert/batch", response_model=BatchResponse)
async def convert_batch(files: list[UploadFile] = File(...)):
    """
    Convertit plusieurs fichiers en une seule requête.

    - **files**: Liste de fichiers (PPTX, PPT, ou PDF)

    Retourne les résultats pour chaque fichier.
    """
    logger.info(f"📦 Conversion batch: {len(files)} fichiers")

    results = []
    success_count = 0
    error_count = 0

    for file in files:
        filename = file.filename or "unnamed"

        try:
            # Déterminer le type de fichier
            if filename.lower().endswith('.pdf'):
                # Conversion PDF
                result = await convert_pdf(file)
            elif filename.lower().endswith(('.pptx', '.ppt')):
                # Conversion PPTX
                result = await convert_pptx(file)
            else:
                results.append(BatchItemResult(
                    filename=filename,
                    status="error",
                    error="Type de fichier non supporté"
                ))
                error_count += 1
                continue

            if result.status == "success":
                results.append(BatchItemResult(
                    filename=filename,
                    status="success",
                    slides=result.slides
                ))
                success_count += 1
            else:
                results.append(BatchItemResult(
                    filename=filename,
                    status="error",
                    error=result.error
                ))
                error_count += 1

        except Exception as e:
            logger.error(f"❌ Erreur batch pour {filename}: {str(e)}")
            results.append(BatchItemResult(
                filename=filename,
                status="error",
                error=str(e)
            ))
            error_count += 1

    logger.info(f"📦 Batch terminé: {success_count} succès, {error_count} erreurs")

    return BatchResponse(
        status="completed",
        results=results,
        total_count=len(files),
        success_count=success_count,
        error_count=error_count
    )


@app.post("/validate")
async def validate_file(file: UploadFile = File(...)):
    """
    Valide un fichier sans le convertir.

    Retourne les métadonnées du fichier.
    """
    filename = file.filename or "unnamed"

    # Vérifier l'extension
    allowed_extensions = ['.pptx', '.ppt', '.pdf']
    ext = os.path.splitext(filename)[1].lower()

    if ext not in allowed_extensions:
        raise HTTPException(
            status_code=400,
            detail=f"Extension non supportée: {ext}. Extensions acceptées: {', '.join(allowed_extensions)}"
        )

    # Lire et vérifier la taille
    content = await file.read()
    file_size = len(content)
    file_size_mb = round(file_size / (1024 * 1024), 2)

    # Limite de 50MB
    if file_size > 50 * 1024 * 1024:
        raise HTTPException(
            status_code=400,
            detail=f"Fichier trop volumineux (max 50MB). Taille actuelle: {file_size_mb}MB"
        )

    return {
        "status": "valid",
        "filename": filename,
        "extension": ext,
        "size": file_size,
        "size_mb": file_size_mb,
        "supported": True
    }


@app.get("/status/{job_id}", response_model=JobStatus)
async def get_job_status(job_id: str):
    """
    Récupère le statut d'un job de conversion.

    - **job_id**: Identifiant unique du job
    """
    jobs = getattr(app.state, 'jobs', {})

    if job_id not in jobs:
        raise HTTPException(
            status_code=404,
            detail=f"Job non trouvé: {job_id}"
        )

    return jobs[job_id]


# ============================================================================
# Point d'entrée
# ============================================================================

if __name__ == "__main__":
    import uvicorn

    uvicorn.run(
        app,
        host="0.0.0.0",
        port=8000,
        log_level="info"
    )
